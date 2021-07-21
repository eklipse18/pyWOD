from pptx import Presentation


def main(word: str):
    # args = ()
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "test"
    subtitle.text = "1234"

    # definition_slide = prs.slides.add_slide()

    prs.save(f"{word.replace(' ', '_')}.pptx")
